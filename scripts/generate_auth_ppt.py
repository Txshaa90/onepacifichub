from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:  # pragma: no cover
    Image = None
    ImageDraw = None
    ImageFont = None


ROOT = Path(__file__).resolve().parents[1]
IMAGE_DIR = ROOT / "Onepacifichub and Supabase Images"
OUTPUT = ROOT / "OnePacificHub_Authentication_System_Supabase_Integration.pptx"
TEMP_DIR = ROOT / "scripts" / "_ppt_code_images"


SLIDE_W = 13.333
SLIDE_H = 7.5

COLORS = {
    "navy": RGBColor(19, 36, 71),
    "blue": RGBColor(25, 94, 210),
    "cyan": RGBColor(24, 167, 201),
    "ink": RGBColor(30, 41, 59),
    "muted": RGBColor(100, 116, 139),
    "line": RGBColor(203, 213, 225),
    "bg": RGBColor(246, 249, 252),
    "card": RGBColor(255, 255, 255),
    "soft": RGBColor(235, 244, 255),
    "soft2": RGBColor(236, 253, 245),
    "warn": RGBColor(255, 247, 237),
}


def inches(value):
    return Inches(value)


def rgb(name):
    return COLORS[name]


def read_excerpt(rel_path, start, end):
    path = ROOT / rel_path
    lines = path.read_text(encoding="utf-8").splitlines()
    body = "\n".join(lines[start - 1 : end])
    return f"{rel_path} ({start}-{end})\n{body}"


def add_textbox(slide, x, y, w, h, text="", size=18, color="ink", bold=False, font="Aptos"):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = inches(0.08)
    tf.margin_right = inches(0.08)
    tf.margin_top = inches(0.05)
    tf.margin_bottom = inches(0.05)
    if text:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color)
    return box


def set_slide_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("bg")


def add_header(slide, title, subtitle=None):
    add_textbox(slide, 0.45, 0.22, 10.6, 0.45, title, size=24, color="navy", bold=True)
    if subtitle:
        add_textbox(slide, 0.45, 0.67, 11.8, 0.32, subtitle, size=10.5, color="muted")
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, inches(0.45), inches(0.98), inches(12.4), inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb("line")
    line.line.fill.background()


def add_slide_number(slide, number):
    return


def add_card(slide, x, y, w, h, fill="card", line="line", radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius_shape, inches(x), inches(y), inches(w), inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    return shape


def add_bullets(slide, x, y, w, h, heading, bullets, fill="card"):
    add_card(slide, x, y, w, h, fill=fill)
    add_textbox(slide, x + 0.16, y + 0.1, w - 0.32, 0.28, heading, size=13, color="navy", bold=True)
    box = slide.shapes.add_textbox(inches(x + 0.16), inches(y + 0.42), inches(w - 0.32), inches(h - 0.52))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = "Aptos"
        p.font.size = Pt(11)
        p.font.color.rgb = rgb("ink")
        p.level = 0
        p.bullet = True


def add_code_box(slide, x, y, w, h, heading, code):
    add_card(slide, x, y, w, h, fill="card")
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inches(x + 0.15), inches(y + 0.13), inches(1.4), inches(0.28)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = rgb("soft")
    tag.line.fill.background()
    tag_tf = tag.text_frame
    tag_tf.paragraphs[0].text = heading
    tag_tf.paragraphs[0].font.name = "Aptos"
    tag_tf.paragraphs[0].font.bold = True
    tag_tf.paragraphs[0].font.size = Pt(10)
    tag_tf.paragraphs[0].font.color.rgb = rgb("blue")

    box = slide.shapes.add_textbox(inches(x + 0.16), inches(y + 0.47), inches(w - 0.32), inches(h - 0.58))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = inches(0.02)
    tf.margin_right = inches(0.02)
    tf.margin_top = inches(0.02)
    tf.margin_bottom = inches(0.02)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code
    r.font.name = "Consolas"
    r.font.size = Pt(8.0)
    r.font.color.rgb = rgb("ink")
    try:
        tf.fit_text(font_family="Consolas", max_size=8)
    except Exception:
        pass


def add_mini_code_card(slide, x, y, w, h, heading, path_label, lines):
    add_card(slide, x, y, w, h, fill="card")
    add_textbox(slide, x + 0.16, y + 0.12, w - 0.32, 0.24, heading, size=12, color="blue", bold=True)
    add_textbox(slide, x + 0.16, y + 0.35, w - 0.32, 0.2, path_label, size=8.5, color="muted")

    box = slide.shapes.add_textbox(inches(x + 0.16), inches(y + 0.62), inches(w - 0.32), inches(h - 0.74))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "\n".join(lines)
    r.font.name = "Consolas"
    r.font.size = Pt(8)
    r.font.color.rgb = rgb("ink")
    try:
        tf.fit_text(font_family="Consolas", max_size=8)
    except Exception:
        pass


def add_image(slide, image_name, x, y, w, h):
    path = IMAGE_DIR / image_name
    if not path.exists():
        return
    if Image is None:
        slide.shapes.add_picture(str(path), inches(x), inches(y), width=inches(w), height=inches(h))
        return

    with Image.open(path) as img:
        ratio = img.width / img.height

    box_ratio = w / h
    if ratio > box_ratio:
        img_w = w
        img_h = w / ratio
    else:
        img_h = h
        img_w = h * ratio

    left = x + (w - img_w) / 2
    top = y + (h - img_h) / 2
    add_card(slide, x, y, w, h, fill="card")
    slide.shapes.add_picture(str(path), inches(left), inches(top), width=inches(img_w), height=inches(img_h))


def get_code_font(size=24):
    if ImageFont is None:
        return None
    candidates = [
        "C:/Windows/Fonts/consola.ttf",
        "C:/Windows/Fonts/consolab.ttf",
        "C:/Windows/Fonts/cour.ttf",
        "C:/Windows/Fonts/lucon.ttf",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def split_code_chunks(code, max_lines=24):
    lines = code.splitlines()
    chunks = []
    for i in range(0, len(lines), max_lines):
        chunk = "\n".join(lines[i : i + max_lines]).rstrip()
        if chunk:
            chunks.append(chunk)
    return chunks or [code]


def make_code_image(name, code, subtitle=""):
    TEMP_DIR.mkdir(parents=True, exist_ok=True)
    out = TEMP_DIR / f"{name}.png"
    if Image is None or ImageDraw is None or ImageFont is None:
        return None

    width = 1800
    height = 980
    img = Image.new("RGB", (width, height), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    draw.rounded_rectangle((20, 20, width - 20, height - 20), radius=48, fill=(255, 255, 255), outline=(203, 213, 225), width=3)
    draw.rounded_rectangle((60, 54, 320, 108), radius=16, fill=(235, 244, 255), outline=None)

    title_font = get_code_font(24)
    code_font = get_code_font(22)
    meta_font = get_code_font(18)

    draw.text((85, 67), subtitle or "Code", fill=(25, 94, 210), font=title_font)
    y = 138
    for line in code.splitlines():
        draw.text((78, y), line, fill=(30, 41, 59), font=code_font)
        y += 30

    img.save(out)
    return out


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)

    ribbon = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inches(0.58), inches(0.72), inches(2.65), inches(0.42)
    )
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = rgb("soft")
    ribbon.line.fill.background()
    ribbon.text_frame.paragraphs[0].text = "Presentation Deck"
    ribbon.text_frame.paragraphs[0].font.name = "Aptos"
    ribbon.text_frame.paragraphs[0].font.size = Pt(12)
    ribbon.text_frame.paragraphs[0].font.bold = True
    ribbon.text_frame.paragraphs[0].font.color.rgb = rgb("blue")

    add_textbox(slide, 0.62, 1.35, 9.6, 1.0, "OnePacificHub Authentication System", size=30, color="navy", bold=True)
    add_textbox(slide, 0.62, 2.15, 7.8, 0.42, "Supabase Integration", size=20, color="cyan", bold=True)
    add_textbox(
        slide,
        0.62,
        2.78,
        7.8,
        0.7,
        "A clear walkthrough of how login, registration, OAuth, session handling, and password reset are implemented in the website.",
        size=12,
        color="ink",
    )

    add_bullets(
        slide,
        0.62,
        3.75,
        5.3,
        1.45,
        "Presentation Focus",
        [
            "Frontend UI to Supabase Auth flow",
            "Supabase dashboard configuration",
            "Actual code integration points in the project",
        ],
        fill="card",
    )

    add_bullets(
        slide,
        6.35,
        1.25,
        6.0,
        2.45,
        "Main Topics",
        [
            "Login and registration",
            "OAuth providers",
            "Password reset with OTP",
            "Supabase dashboard setup and code flow",
        ],
        fill="card",
    )

    add_slide_number(slide, 1)


def add_image_only_slide(prs, number, title, subtitle, image_name, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, title, subtitle)
    add_image(slide, image_name, 0.45, 1.18, 7.1, 5.7)
    add_bullets(slide, 7.8, 1.18, 5.05, 4.75, "Key Points", bullets, fill="card")
    add_slide_number(slide, number)


def add_text_only_slide(prs, number, title, subtitle, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, title, subtitle)
    add_bullets(slide, 0.9, 1.35, 11.5, 5.25, "Summary", bullets, fill="card")
    add_slide_number(slide, number)


def add_code_detail_slides(prs, start_number, title, subtitle, code_heading, code):
    chunks = split_code_chunks(code, max_lines=24)
    number = start_number
    for idx, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide)
        part_title = title if len(chunks) == 1 else f"{title} ({idx}/{len(chunks)})"
        add_header(slide, part_title, subtitle)
        image_path = make_code_image(f"{number}_{idx}", chunk, code_heading)
        if image_path:
            slide.shapes.add_picture(str(image_path), inches(0.72), inches(1.28), width=inches(11.9), height=inches(5.9))
        else:
            add_code_box(slide, 0.7, 1.35, 11.95, 5.7, code_heading, chunk)
        add_slide_number(slide, number)
        number += 1
    return number


def add_flow_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Full Auth Flow", "Frontend to Supabase session lifecycle")

    labels = [
        ("User", 0.72),
        ("Login UI", 2.55),
        ("Auth Service", 4.58),
        ("Supabase", 7.0),
        ("Session", 9.3),
        ("Protected Page", 11.0),
    ]

    for label, x in labels:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inches(x), inches(2.1), inches(1.55), inches(0.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb("card")
        shape.line.color.rgb = rgb("blue")
        tf = shape.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.name = "Aptos"
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = rgb("navy")

    for x in [2.0, 4.0, 6.45, 8.75, 10.55]:
        arrow = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON, inches(x), inches(2.35), inches(0.45), inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = rgb("cyan")
        arrow.line.fill.background()

    add_bullets(
        slide,
        0.72,
        3.55,
        5.8,
        2.1,
        "Flow Explanation",
        [
            "User interacts with login, signup, OAuth, or reset screens.",
            "Service layer calls Supabase Auth or Edge Functions.",
            "Supabase returns a session that the app uses to protect routes.",
        ],
        fill="card",
    )

    pseudo = (
        "User -> LoginPage/RegisterPage\n"
        "     -> supabaseAuthService\n"
        "     -> Supabase Auth / Edge Functions\n"
        "     -> session token\n"
        "     -> AuthContext + ProtectedRoute\n"
        "     -> Account page"
    )
    add_code_box(slide, 6.75, 3.55, 6.1, 2.1, "Flow Diagram", pseudo)
    add_slide_number(slide, number)


def add_integration_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Code Integration", "Three files that connect the whole auth stack")

    add_bullets(
        slide,
        0.7,
        1.35,
        3.8,
        3.2,
        "supabase.js",
        [
            "Initializes the browser client.",
            "Reads Supabase URL and anon key from environment variables.",
            "Configures auth persistence and session handling.",
        ],
        fill="card",
    )
    add_bullets(
        slide,
        4.78,
        1.35,
        3.8,
        3.2,
        "AuthContext.jsx",
        [
            "Watches auth state changes.",
            "Stores the logged-in user in React state.",
            "Protects routes and syncs session data.",
        ],
        fill="card",
    )
    add_bullets(
        slide,
        8.86,
        1.35,
        3.8,
        3.2,
        "supabaseAuthService.js",
        [
            "Wraps login, signup, OAuth, logout, and session checks.",
            "Keeps Supabase calls out of page components.",
            "Returns simplified user/session objects to the app.",
        ],
        fill="card",
    )
    add_bullets(
        slide,
        0.95,
        4.95,
        11.7,
        1.25,
        "How They Work Together",
        [
            "The client connects to Supabase, the service performs authentication requests, and AuthContext keeps the session active across the website.",
        ],
        fill="soft",
    )

    add_slide_number(slide, number)


def add_limitation_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Current Limitation", "Supabase is used for authentication only in the current frontend")

    add_bullets(
        slide,
        0.9,
        1.35,
        11.5,
        4.9,
        "What This Means",
        [
            "Login, signup, OAuth, session handling, and reset flow are on Supabase.",
            "Profile update and password-change helpers still call a mock local-storage service.",
            "Products, cart, and broader app data are not being queried from Supabase here.",
            "So the current setup is production-style auth, but not yet a full Supabase backend.",
        ],
        fill="warn",
    )
    add_slide_number(slide, number)


def add_conclusion_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Conclusion", "Why this integration is effective for the project")

    add_bullets(
        slide,
        0.62,
        1.28,
        3.82,
        2.0,
        "Strengths",
        [
            "Secure authentication is delegated to Supabase.",
            "Frontend stays focused on UX and route protection.",
            "OAuth and password reset are already integrated.",
        ],
        fill="card",
    )
    add_bullets(
        slide,
        4.76,
        1.28,
        3.82,
        2.0,
        "Architecture",
        [
            "UI pages collect input.",
            "Services call Supabase APIs.",
            "AuthContext stores the active user session.",
        ],
        fill="card",
    )
    add_bullets(
        slide,
        8.9,
        1.28,
        3.82,
        2.0,
        "Next Step",
        [
            "Move profile and app data to Supabase if full backend consolidation is desired.",
            "Add database tables and replace mock services.",
        ],
        fill="soft2",
    )

    add_code_box(
        slide,
        1.6,
        3.72,
        10.15,
        1.85,
        "Final Takeaway",
        "Supabase simplifies authentication by handling login, sessions, and security,\nallowing the frontend to focus on user experience and route protection.",
    )
    add_slide_number(slide, number)


def build_presentation():
    prs = Presentation()
    prs.slide_width = inches(SLIDE_W)
    prs.slide_height = inches(SLIDE_H)

    add_title_slide(prs)
    n = 2

    add_text_only_slide(
        prs, n, "What Supabase Does Here", "Supabase is the auth layer, not the full application backend",
        ["Handles login and signup", "Manages sessions and token refresh", "Supports Google and Facebook OAuth", "Supports password reset through Edge Functions."]
    )
    n += 1
    n = add_code_detail_slides(prs, n, "Supabase Client Code", "Core browser client setup", "src/lib/supabase.js", read_excerpt("src/lib/supabase.js", 1, 24))

    add_image_only_slide(
        prs, n, "Sign In Page", "Email/password login plus social login entry point", "Sign In.png",
        ["Users can log in with email and password.", "The same screen also exposes Google and Facebook login.", "Successful login creates a Supabase session."]
    )
    n += 1
    n = add_code_detail_slides(prs, n, "Sign In Code", "Email/password authentication logic", "src/services/supabaseAuthService.js", read_excerpt("src/services/supabaseAuthService.js", 4, 22))

    add_image_only_slide(
        prs, n, "Sign Up Page", "User registration form connected directly to Supabase Auth", "Sign Up.png",
        ["The form collects first name, last name, email, and password.", "The signup request is sent to Supabase Auth.", "Verification can be required before first login."]
    )
    n += 1
    n = add_code_detail_slides(prs, n, "Sign Up Code", "Account creation and verification flow", "src/services/supabaseAuthService.js", read_excerpt("src/services/supabaseAuthService.js", 73, 129))

    add_image_only_slide(
        prs, n, "OAuth Login", "Social authentication through Supabase providers", "Sign In.png",
        ["Google and Facebook buttons are available in the login UI.", "Supabase handles the provider redirect and callback.", "OAuth returns the user to the app with a valid session."]
    )
    n += 1
    n = add_code_detail_slides(prs, n, "OAuth Code", "Google and Facebook login logic", "src/services/supabaseAuthService.js", read_excerpt("src/services/supabaseAuthService.js", 34, 70))

    add_image_only_slide(
        prs, n, "Forgot Password", "Password reset starts by requesting a one-time code", "Forgot password.png",
        ["The user enters an email address to start recovery.", "The app does not use the default magic-link reset flow.", "It calls a custom Edge Function to send an OTP."]
    )
    n += 1
    n = add_code_detail_slides(prs, n, "Forgot Password Code", "Reset code request through Edge Functions", "src/services/passwordResetOtp.js", read_excerpt("src/services/passwordResetOtp.js", 20, 49))

    add_image_only_slide(
        prs, n, "OTP Verification", "The reset flow continues with a six-digit code", "Forgot password send code.png",
        ["The user enters the one-time code sent by email.", "The backend validates the OTP and returns a short-lived reset token.", "That token is then used to set the new password."]
    )
    n += 1
    n = add_code_detail_slides(prs, n, "OTP Verification Code", "OTP validation and reset token handling", "src/services/passwordResetOtp.js", read_excerpt("src/services/passwordResetOtp.js", 52, 84))

    add_image_only_slide(
        prs, n, "Email Received", "The code is delivered to the user outside the app UI", "OTP Sent to email.png",
        ["The email contains the verification code used in the reset flow.", "The code expires after a limited time window.", "This keeps recovery time-bound and more secure."]
    )
    n += 1

    add_image_only_slide(
        prs, n, "Logged-in Account Page", "Protected pages become available after authentication succeeds", "Sign in with google profile.png",
        ["After login, the user is redirected to the account dashboard.", "The route is protected by authentication state.", "Session changes are monitored through AuthContext."]
    )
    n += 1
    n = add_code_detail_slides(prs, n, "Session Code", "AuthContext session subscription and sync", "src/context/AuthContext.jsx", read_excerpt("src/context/AuthContext.jsx", 23, 61))

    add_image_only_slide(
        prs, n, "Supabase Users", "Registered users are stored in Supabase Authentication", "Authentication - Users.png",
        ["The dashboard stores each user record in Supabase Auth.", "Key details include email, provider, and user ID.", "Custom metadata such as first and last name can also be attached."]
    )
    n += 1

    add_image_only_slide(
        prs, n, "Auth Providers", "Enabled providers determine which login methods can be used", "Authentication - Providers.png",
        ["Email login is configured in Supabase.", "Google and Facebook can be enabled as OAuth providers.", "Only enabled providers should be exposed in the UI."]
    )
    n += 1
    n = add_code_detail_slides(prs, n, "Provider Button Code", "Frontend buttons for social login", "src/components/SocialOAuthButtons.jsx", read_excerpt("src/components/SocialOAuthButtons.jsx", 20, 58))

    add_image_only_slide(
        prs, n, "URL Configuration", "Redirect URLs connect the dashboard settings to the frontend callback route", "URL Configuration.png",
        ["Supabase needs approved redirect URLs for OAuth and email flows.", "The site origin is built from environment variables.", "This must match the AuthCallback route used by the app."]
    )
    n += 1
    n = add_code_detail_slides(prs, n, "URL Callback Code", "Redirect helpers and callback exchange", "siteUrl.js + AuthCallback.jsx", read_excerpt("src/lib/siteUrl.js", 1, 24) + "\n\n" + read_excerpt("src/pages/AuthCallback.jsx", 54, 80))

    add_image_only_slide(
        prs, n, "Edge Functions", "Custom backend logic is used for the password reset flow", "Edge Functions.png",
        ["The app expects custom functions named send-reset-otp, verify-reset-otp, and reset-password.", "Those functions run behind Supabase Functions.", "They extend Supabase Auth with custom OTP-based recovery."]
    )
    n += 1
    n = add_code_detail_slides(prs, n, "Edge Function Client Code", "Frontend calls to Supabase Functions", "src/services/passwordResetOtp.js", read_excerpt("src/services/passwordResetOtp.js", 1, 84))

    add_flow_slide(prs, n)
    n += 1
    add_integration_slide(prs, n)
    n += 1
    n = add_code_detail_slides(
        prs, n, "Integration Code Detail", "The three main files behind the auth system", "supabase.js + AuthContext + service",
        read_excerpt("src/lib/supabase.js", 1, 24) + "\n\n" + read_excerpt("src/context/AuthContext.jsx", 64, 122) + "\n\n" + read_excerpt("src/services/supabaseAuthService.js", 147, 186)
    )
    add_limitation_slide(prs, n)
    n += 1
    n = add_code_detail_slides(
        prs, n, "Limitation Code Detail", "Mock/local service still used outside auth", "src/services/authService.js",
        read_excerpt("src/services/authService.js", 1, 16) + "\n\n" + read_excerpt("src/services/authService.js", 223, 277)
    )
    add_conclusion_slide(prs, n)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
    print(OUTPUT)
